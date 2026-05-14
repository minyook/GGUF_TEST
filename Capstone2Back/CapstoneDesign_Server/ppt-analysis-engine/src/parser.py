"""PPT 파일에서 원천 데이터(텍스트/노트/도형/폰트 등)를 최대한 추출합니다.

이 모듈은 웹 프레임워크와 독립적으로 동작하도록, 파일 경로 입력 → dict 반환 형태로 구성됩니다.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide

ProgressCallback = Callable[[str], None]


def _strip_text(value: str | None) -> str:
    if not value:
        return ""
    return value.strip()


class PPTParser:
    """python-pptx 기반 발표 자료 파서."""

    def __init__(self, ppt_path: str | Path) -> None:
        self.ppt_path = Path(ppt_path)

    def extract(self, progress_callback: ProgressCallback | None = None) -> dict[str, Any]:
        if not self.ppt_path.is_file():
            raise FileNotFoundError(f"PPT 파일을 찾을 수 없습니다: {self.ppt_path}")

        if progress_callback:
            progress_callback("PPT 파일 로딩 중...")
        prs = Presentation(str(self.ppt_path))
        slides_out: list[dict[str, Any]] = []
        total_slides = len(prs.slides)

        if progress_callback:
            progress_callback(f"슬라이드 읽기 시작 (총 {total_slides}장)")
            progress_callback("읽는 정보: 제목, 본문 텍스트, 이미지 위치, 텍스트 박스, 글자 크기, 발표자 노트")

        for idx, slide in enumerate(prs.slides):
            if progress_callback:
                current = idx + 1
                if current == 1 or current % 5 == 0 or current == total_slides:
                    progress_callback(f"슬라이드 정보 읽는 중... ({current}/{total_slides})")
            slides_out.append(self._extract_slide(idx, slide))

        core = prs.core_properties
        return {
            "source_path": str(self.ppt_path.resolve()),
            "file_name": self.ppt_path.name,
            "slide_count": len(slides_out),
            "slide_size_emu": {"width": int(prs.slide_width), "height": int(prs.slide_height)},
            "metadata": {
                "title": getattr(core, "title", None),
                "subject": getattr(core, "subject", None),
                "author": getattr(core, "author", None),
                "company": getattr(core, "company", None),
                "created": core.created.isoformat() if getattr(core, "created", None) else None,
                "modified": core.modified.isoformat() if getattr(core, "modified", None) else None,
            },
            "slides": slides_out,
        }

    def _extract_slide(self, slide_index: int, slide: Slide) -> dict[str, Any]:
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = _strip_text(slide.shapes.title.text)

        body_chunks: list[str] = []
        notes_text = self._extract_notes_text(slide)
        image_count = 0
        shape_count = 0
        textbox_count = 0
        autoshape_count = 0
        fonts: list[dict[str, Any]] = []
        image_items: list[dict[str, Any]] = []
        text_boxes: list[dict[str, Any]] = []
        shapes_by_type: dict[str, int] = {}
        text_shape_area_emu2 = 0

        title_shape = slide.shapes.title
        for shape in self._iter_shapes(slide.shapes):
            shape_count += 1
            st_name = getattr(shape, "shape_type", None)
            st_key = str(st_name)
            shapes_by_type[st_key] = shapes_by_type.get(st_key, 0) + 1

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image_items.append(self._image_info(shape, image_count))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                autoshape_count += 1

            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                if title_shape is not None and shape is title_shape:
                    continue
                textbox_count += 1
                try:
                    # EMU 단위 면적(대략) = width * height
                    text_shape_area_emu2 += int(shape.width) * int(shape.height)
                except Exception:
                    pass
                for paragraph in shape.text_frame.paragraphs:
                    para_text = _strip_text(paragraph.text)
                    if para_text:
                        body_chunks.append(para_text)
                    for run in paragraph.runs:
                        font_info = self._font_info_from_run(run)
                        if font_info:
                            fonts.append(font_info)
                text_boxes.append(self._text_box_info(shape))

        body_text_merged = "\n".join(body_chunks)

        return {
            "slide_index": slide_index,
            "title": title,
            "body_text": body_text_merged,
            "body_segments": body_chunks,
            "notes_text": notes_text,
            "image_count": image_count,
            "shape_count": shape_count,
            "textbox_count": textbox_count,
            "autoshape_count": autoshape_count,
            "shapes_by_type": shapes_by_type,
            "text_shape_area_emu2": text_shape_area_emu2,
            "image_items": image_items,
            "text_boxes": text_boxes,
            "fonts": fonts,
        }

    def _iter_shapes(self, shapes: Any):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes(shape.shapes)

    def _extract_notes_text(self, slide: Slide) -> str:
        # notes_slide는 존재하더라도 notes_text_frame이 비어있을 수 있음
        try:
            notes_slide = slide.notes_slide
        except Exception:
            return ""
        try:
            tf = notes_slide.notes_text_frame
        except Exception:
            return ""
        try:
            return _strip_text(tf.text)
        except Exception:
            return ""

    def _font_info_from_run(self, run: Any) -> dict[str, Any] | None:
        text = _strip_text(run.text)
        if not text:
            return None

        font = run.font
        size_pt: float | None = None
        if font.size is not None:
            try:
                size_pt = float(font.size.pt)
            except (AttributeError, TypeError):
                size_pt = None

        return {
            "text_sample": text[:200],
            "name": font.name,
            "size_pt": size_pt,
            "bold": bool(font.bold) if font.bold is not None else None,
            "italic": bool(font.italic) if font.italic is not None else None,
        }

    def _shape_bounds(self, shape: Any) -> dict[str, int]:
        return {
            "left": int(getattr(shape, "left", 0) or 0),
            "top": int(getattr(shape, "top", 0) or 0),
            "width": int(getattr(shape, "width", 0) or 0),
            "height": int(getattr(shape, "height", 0) or 0),
        }

    def _image_info(self, shape: Any, fallback_index: int) -> dict[str, Any]:
        name = getattr(shape, "name", None) or f"image_{fallback_index}"
        info = {"name": name, **self._shape_bounds(shape)}
        return info

    def _text_box_info(self, shape: Any) -> dict[str, Any]:
        txt = ""
        try:
            txt = _strip_text(shape.text_frame.text)
        except Exception:
            txt = ""
        return {"text": txt[:300], **self._shape_bounds(shape)}


def parse_ppt_file(
    ppt_path: str | Path,
    progress_callback: ProgressCallback | None = None,
) -> dict[str, Any]:
    """파일 경로를 받아 파싱 결과 딕셔너리를 반환합니다."""
    return PPTParser(ppt_path).extract(progress_callback=progress_callback)
